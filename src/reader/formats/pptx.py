from html import escape
from pathlib import Path

from pptx import Presentation

from reader.preview.result import PreviewResult


def to_html(path: Path) -> PreviewResult:
    prs = Presentation(str(path))
    sections: list[str] = []
    for index, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = "".join(run.text for run in para.runs) or para.text
                if t:
                    texts.append(f"<p>{escape(t)}</p>")
            table = None
            try:
                if shape.has_table:
                    table = shape.table
            except AttributeError:
                pass
            if table is not None:
                rows = ["<table>"]
                for row in table.rows:
                    cells = "".join(f"<td>{escape(cell.text)}</td>" for cell in row.cells)
                    rows.append(f"<tr>{cells}</tr>")
                rows.append("</table>")
                texts.append("".join(rows))
        sections.append(
            f'<section class="slide" id="slide-{index}"><h2>Slide {index}</h2>{"".join(texts)}</section>'
        )
    body = "\n".join(sections)
    html = (
        "<!DOCTYPE html><html><head><meta charset='utf-8'>"
        "<style>.slide{border:1px solid #ddd;margin:12px;padding:12px}</style>"
        f"</head><body>{body}</body></html>"
    )
    return PreviewResult(html=html, status_label="内置预览", kind="html")
