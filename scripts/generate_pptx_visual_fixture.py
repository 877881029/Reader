import argparse
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
import re
from zipfile import ZIP_STORED, ZipFile, ZipInfo

from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


DEFAULT_OUTPUT = (
    Path(__file__).resolve().parents[1]
    / "tests"
    / "fixtures"
    / "pptx"
    / "visual-elements.pptx"
)
_ZIP_TIMESTAMP = (2000, 1, 1, 0, 0, 0)
_CORE_TIME = re.compile(
    rb"(<dcterms:(?:created|modified)[^>]*>).*?"
    rb"(</dcterms:(?:created|modified)>)"
)


def _build_presentation() -> Presentation:
    presentation = Presentation()
    fixed_time = datetime(2000, 1, 1, tzinfo=timezone.utc)
    presentation.core_properties.created = fixed_time
    presentation.core_properties.modified = fixed_time

    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(20, 48, 90)
    slide.shapes.title.text = "Inherited title"
    image = BytesIO()
    Image.new("RGB", (320, 180), (37, 99, 235)).save(image, "PNG")
    image.seek(0)
    slide.shapes.add_picture(image, Inches(7), Inches(2), width=Inches(4))

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(
        3, 3, Inches(1), Inches(1), Inches(10), Inches(3)
    ).table
    for row_index, row in enumerate(
        (("Metric", "Q1", "Q2"), ("A", "10", "14"), ("B", "8", "13"))
    ):
        for column_index, value in enumerate(row):
            table.cell(row_index, column_index).text = value

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    chart_data = ChartData()
    chart_data.categories = ["North", "South", "West"]
    chart_data.add_series("Revenue", (12, 18, 15))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(10),
        Inches(5),
        chart_data,
    )

    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    box.text = "Missing font continues"
    box.text_frame.paragraphs[0].runs[0].font.name = "ReaderMissingFontZZ"
    return presentation


def _normalize_zip(source_bytes: bytes) -> bytes:
    target = BytesIO()
    with ZipFile(BytesIO(source_bytes), "r") as archive, ZipFile(
        target, "w", compression=ZIP_STORED
    ) as output:
        for name in sorted(archive.namelist()):
            payload = archive.read(name)
            if name == "docProps/core.xml":
                payload = _CORE_TIME.sub(
                    rb"\g<1>2000-01-01T00:00:00Z\g<2>", payload
                )
            elif name.lower().endswith((".xlsx", ".xlsm")):
                payload = _normalize_zip(payload)
            info = ZipInfo(name, date_time=_ZIP_TIMESTAMP)
            info.compress_type = ZIP_STORED
            info.create_system = 0
            info.external_attr = 0
            output.writestr(info, payload)
    return target.getvalue()


def _deterministic_package(presentation: Presentation) -> bytes:
    source = BytesIO()
    presentation.save(source)
    return _normalize_zip(source.getvalue())


def generate(output: Path = DEFAULT_OUTPUT) -> Path:
    output = output.resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_bytes(_deterministic_package(_build_presentation()))
    return output


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Generate the deterministic visual PPTX integration fixture."
    )
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    args = parser.parse_args()
    generate(args.output)


if __name__ == "__main__":
    main()
