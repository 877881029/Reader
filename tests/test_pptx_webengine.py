from __future__ import annotations

import hashlib
import itertools
import json
import subprocess
import sys
from pathlib import Path

import pytest
import shiboken6
from PySide6.QtCore import QEvent, QEventLoop, QLibraryInfo, QTimer, QUrl
from PySide6.QtWebEngineCore import QWebEngineSettings
from PySide6.QtWebEngineWidgets import QWebEngineView

from reader.preview.pptx_view import PptxVisualView
from reader.preview.result import PreviewResult
from reader.resources import resource_path
from reader.shell.window import MainWindow


ROOT = Path(__file__).resolve().parents[1]
FIXTURE = ROOT / "tests" / "fixtures" / "pptx" / "visual-elements.pptx"
GENERATOR = ROOT / "scripts" / "generate_pptx_visual_fixture.py"
TRACKER_URL = "https://example.invalid/tracker.png"
OUTBOUND_URLS = (
    "http://example.invalid/pixel.png",
    TRACKER_URL,
    "ws://example.invalid/socket",
    "wss://example.invalid/socket",
)
_js_ids = itertools.count()

qt_root = Path(
    QLibraryInfo.path(QLibraryInfo.LibraryPath.LibraryExecutablesPath)
)
WEBENGINE_MISSING = (
    not resource_path("assets", "pptx-viewer", "index.html").is_file()
    or not (qt_root / "QtWebEngineProcess.exe").is_file()
)

pytestmark = [
    pytest.mark.webengine,
    pytest.mark.skipif(
        WEBENGINE_MISSING,
        reason="committed WebEngine bundle unavailable",
    ),
]


def _result() -> PreviewResult:
    return PreviewResult(
        html="",
        fallback_html="<p>safe fallback</p>",
        status_label="内置预览",
        kind="pptx",
    )


def run_js(qtbot, view: QWebEngineView, script: str):
    """Run JavaScript with a unique callback owner and a condition wait."""
    key = f"js-{next(_js_ids)}"
    completed: dict[str, object] = {}
    view.page().runJavaScript(
        script,
        lambda value, callback_key=key: completed.__setitem__(
            callback_key, value
        ),
    )
    qtbot.waitUntil(lambda: key in completed, timeout=10_000)
    return completed[key]


def run_json(qtbot, view: QWebEngineView, expression: str):
    return json.loads(
        run_js(qtbot, view, f"JSON.stringify({expression})")
    )


@pytest.fixture(scope="session", autouse=True)
def webengine_process_probe(qapp):
    if WEBENGINE_MISSING:
        yield
        return

    view = QWebEngineView()
    loop = QEventLoop()
    outcome: list[bool] = []
    deadline = QTimer()
    deadline.setSingleShot(True)
    deadline.timeout.connect(loop.quit)

    def loaded(succeeded: bool) -> None:
        outcome.append(succeeded)
        loop.quit()

    view.loadFinished.connect(loaded)
    deadline.start(5_000)
    view.load(QUrl("data:text/html,<title>reader-webengine-probe</title>ok"))
    loop.exec()
    deadline.stop()
    view.stop()
    view.close()
    view.deleteLater()
    qapp.sendPostedEvents(None, QEvent.Type.DeferredDelete)
    qapp.processEvents()

    if outcome != [True]:
        pytest.skip("QtWebEngine process cannot start")
    yield


def _start_real_view(qtbot, *, fail_slide: int | None = None):
    view = PptxVisualView(_result(), FIXTURE, test_fail_slide=fail_slide)
    qtbot.addWidget(view)
    view.resize(1200, 800)
    view.show()
    ready: list[int] = []
    changed: list[int] = []
    failures: list[str] = []
    view.ready.connect(ready.append)
    view.slide_changed.connect(changed.append)
    view.render_failed.connect(failures.append)
    view.start()
    qtbot.waitUntil(lambda: ready == [4], timeout=20_000)
    return view, changed, failures


def test_fixture_generator_is_byte_deterministic_and_contains_real_elements(
    tmp_path,
):
    first = tmp_path / "first.pptx"
    second = tmp_path / "second.pptx"
    for output in (first, second):
        subprocess.run(
            [sys.executable, str(GENERATOR), "--output", str(output)],
            cwd=ROOT,
            check=True,
        )

    assert hashlib.sha256(first.read_bytes()).digest() == hashlib.sha256(
        second.read_bytes()
    ).digest()
    assert first.read_bytes() == FIXTURE.read_bytes()

    from pptx import Presentation

    presentation = Presentation(first)
    assert len(presentation.slides) == 4
    assert presentation.slides[0].background.fill.fore_color.rgb is not None
    assert any(
        shape.shape_type is not None and hasattr(shape, "image")
        for shape in presentation.slides[0].shapes
    )
    assert any(shape.has_table for shape in presentation.slides[1].shapes)
    assert any(shape.has_chart for shape in presentation.slides[2].shapes)
    assert (
        presentation.slides[3]
        .shapes[0]
        .text_frame.paragraphs[0]
        .runs[0]
        .font.name
        == "ReaderMissingFontZZ"
    )


def test_real_view_renders_fidelity_navigation_zoom_fit_and_blocks_network(
    qtbot,
):
    view, changed, failures = _start_real_view(qtbot)
    interceptor = view.interceptor
    assert interceptor is not None

    first_slide = run_json(
        qtbot,
        view,
        """
        (() => ({
          count: document.querySelectorAll("[data-slide-index]").length,
          active: document.querySelector(".viewer-shell__thumb.is-active")
            ?.getAttribute("data-slide-index"),
          image: Boolean(document.querySelector(".viewer-shell__host svg image")),
          types: document.querySelector("#app")?.getAttribute("data-element-types"),
        }))()
        """,
    )
    assert first_slide["count"] == 4
    assert first_slide["active"] == "0"
    assert first_slide["image"] is True
    element_types = json.loads(first_slide["types"])
    assert "image" in element_types[0].split(",")
    assert "table" in element_types[1].split(",")
    assert "chart" in element_types[2].split(",")
    assert "text" in element_types[3].split(",")

    assert run_json(
        qtbot,
        view,
        """
        (() => {
          document.querySelector('[data-slide-index="1"]').click();
          return {
            page: document.querySelector(".viewer-shell__page").textContent,
            table: Boolean(document.querySelector(
              ".viewer-shell__host foreignObject table"
            )),
          };
        })()
        """,
    ) == {"page": "2 / 4", "table": True}
    qtbot.waitUntil(lambda: changed[-1:] == [1])

    chart = run_json(
        qtbot,
        view,
        """
        (() => {
          document.querySelector('[data-slide-index="2"]').click();
          return {
            structures: document.querySelectorAll(
              ".viewer-shell__host svg path, .viewer-shell__host svg rect"
            ).length,
            chart: JSON.parse(document.querySelector("#app").dataset.elementTypes)
              [2].split(",").includes("chart"),
          };
        })()
        """,
    )
    assert chart["chart"] is True
    assert chart["structures"] > 3

    key_expectations = (
        ("Home", "1 / 4", 0),
        ("ArrowRight", "2 / 4", 1),
        ("PageDown", "3 / 4", 2),
        ("End", "4 / 4", 3),
        ("ArrowLeft", "3 / 4", 2),
        ("PageUp", "2 / 4", 1),
    )
    for key, page, index in key_expectations:
        assert (
            run_js(
                qtbot,
                view,
                f"""
                (() => {{
                  const root = document.querySelector("#app");
                  root.focus();
                  root.dispatchEvent(new KeyboardEvent("keydown", {{
                    key: {key!r}, bubbles: true
                  }}));
                  return document.querySelector(
                    ".viewer-shell__page"
                  ).textContent;
                }})()
                """,
            )
            == page
        )
        qtbot.waitUntil(lambda expected=index: changed[-1:] == [expected])

    assert (
        run_js(
            qtbot,
            view,
            """
            (() => {
              document.querySelector('[data-action="zoom-in"]').click();
              return document.querySelector(".viewer-shell__zoom").textContent;
            })()
            """,
        )
        != "100%"
    )

    run_js(
        qtbot,
        view,
        """document.querySelector('[data-slide-index="3"]').click(); true""",
    )
    qtbot.waitUntil(lambda: changed[-1:] == [3])
    assert run_json(
        qtbot,
        view,
        """Boolean(document.querySelector(".viewer-shell__host svg"))""",
    )

    before = run_json(
        qtbot,
        view,
        """
        (() => {
          document.querySelector('[data-action="fit"]').click();
          return {
            zoom: document.querySelector(".viewer-shell__zoom").textContent,
            width: document.querySelector(".viewer-shell__stage").clientWidth,
          };
        })()
        """,
    )
    view.resize(1500, 900)
    qtbot.waitUntil(
        lambda: run_js(
            qtbot,
            view,
            "document.querySelector('.viewer-shell__stage').clientWidth",
        )
        != before["width"],
        timeout=10_000,
    )
    after_zoom = run_js(
        qtbot,
        view,
        """
        (() => {
          document.querySelector('[data-action="fit"]').click();
          return document.querySelector(".viewer-shell__zoom").textContent;
        })()
        """,
    )
    assert after_zoom != before["zoom"]

    # Temporarily bypass the first local-content policy layer so the real
    # Chromium requests reach and exercise the interceptor layer.
    view.page().settings().setAttribute(
        QWebEngineSettings.WebAttribute.LocalContentCanAccessRemoteUrls, True
    )
    run_js(
        qtbot,
        view,
        """
        (() => {
          const urls = [
            "http://example.invalid/pixel.png",
            "https://example.invalid/tracker.png",
          ];
          for (const source of urls) {
            const image = document.createElement("img");
            image.dataset.outboundProbe = source;
            image.src = source;
            document.body.append(image);
          }
          for (const source of [
            "ws://example.invalid/socket",
            "wss://example.invalid/socket",
          ]) {
            try { new WebSocket(source); } catch {}
          }
          return true;
        })()
        """,
    )
    qtbot.waitUntil(
        lambda: set(OUTBOUND_URLS).issubset(
            set(interceptor.blocked_urls())
        ),
        timeout=10_000,
    )
    blocked = interceptor.blocked_urls()
    assert blocked.count(TRACKER_URL) == 1
    assert set(OUTBOUND_URLS).issubset(set(blocked))
    assert run_json(
        qtbot,
        view,
        """
        (() => {
          const image = [...document.querySelectorAll("[data-outbound-probe]")]
            .find(node => node.dataset.outboundProbe.includes("tracker.png"));
          return image.complete && image.naturalWidth === 0;
        })()
        """,
    )
    assert failures == []
    view.shutdown()


def test_single_slide_failure_stays_local_and_other_slides_continue(qtbot):
    view, changed, failures = _start_real_view(qtbot, fail_slide=1)

    assert "testFailSlide" not in view.viewer_url.query()
    assert run_json(
        qtbot,
        view,
        """
        (() => {
          document.querySelector('[data-slide-index="1"]').click();
          return {
            error: Boolean(document.querySelector(
              ".viewer-shell__slide-error"
            )),
            text: document.querySelector(
              ".viewer-shell__slide-error"
            )?.textContent,
          };
        })()
        """,
    ) == {"error": True, "text": "第 2 页无法渲染"}
    qtbot.waitUntil(lambda: changed[-1:] == [1])

    assert run_json(
        qtbot,
        view,
        """
        (() => {
          document.querySelector('[data-slide-index="2"]').click();
          return {
            svg: Boolean(document.querySelector(".viewer-shell__host svg")),
            error: Boolean(document.querySelector(
              ".viewer-shell__slide-error"
            )),
          };
        })()
        """,
    ) == {"svg": True, "error": False}
    qtbot.waitUntil(lambda: changed[-1:] == [2])
    assert failures == []
    assert not view.is_fallback
    view.shutdown()


class _OfflineOffice:
    def available_for(self, _suffix: str) -> bool:
        return False


def test_main_window_default_factory_loads_and_releases_webengine_profile(
    qtbot,
):
    window = MainWindow(
        preview_fn=lambda *_args, **_kwargs: _result(),
        office=_OfflineOffice(),
    )
    qtbot.addWidget(window)
    window.show()
    window.open_paths([str(FIXTURE)])

    qtbot.waitUntil(
        lambda: bool(window._documents)
        and next(iter(window._documents.values())).visual_slide_count == 4,
        timeout=20_000,
    )
    document = next(iter(window._documents.values()))
    layout = document.page.layout()
    assert layout is not None
    content = layout.itemAt(0).widget()
    assert isinstance(content, PptxVisualView)
    profile = content.profile
    assert profile is not None and profile.isOffTheRecord()

    window.close()
    qtbot.waitUntil(lambda: not shiboken6.isValid(profile), timeout=10_000)
    assert content.profile is None
    assert not window._documents
