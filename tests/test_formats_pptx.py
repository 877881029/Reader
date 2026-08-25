from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from reader.formats import pptx as fmt_pptx
from reader.formats.pptx import to_html, to_visual


def test_pptx_emits_one_section_per_slide(tmp_path: Path):
    path = tmp_path / "s.pptx"
    prs = Presentation()
    layout = prs.slide_layouts[1]
    s1 = prs.slides.add_slide(layout)
    s1.shapes.title.text = "Slide One"
    s1.placeholders[1].text = "Alpha"
    s2 = prs.slides.add_slide(layout)
    s2.shapes.title.text = "Slide Two"
    s2.placeholders[1].text = "Beta"
    prs.save(path)
    result = to_html(path)
    assert result.html.count('class="slide"') == 2
    assert "Slide One" in result.html
    assert "Alpha" in result.html
    assert "Slide Two" in result.html
    assert "Beta" in result.html
    assert result.status_label == "内置预览"
    assert result.kind == "html"
    assert result.error is None


def test_pptx_renders_table_on_slide(tmp_path: Path):
    path = tmp_path / "table.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table
    table.cell(0, 0).text = "H1"
    table.cell(0, 1).text = "H2"
    table.cell(1, 0).text = "a"
    table.cell(1, 1).text = "b"
    prs.save(path)
    result = to_html(path)
    assert result.status_label == "内置预览"
    assert result.kind == "html"
    assert result.error is None
    assert "H1" in result.html
    assert "H2" in result.html
    assert "a" in result.html
    assert "b" in result.html
    assert "<table" in result.html.lower()


def test_pptx_escapes_dangerous_characters(tmp_path: Path):
    path = tmp_path / "escape.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = '<script>alert(1)</script> & "quotes"'
    table = slide.shapes.add_table(1, 1, Inches(1), Inches(2), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "<img onerror=alert(1)>"
    prs.save(path)
    result = to_html(path)
    html = result.html
    assert result.error is None
    assert "<script>" not in html
    assert "&lt;script&gt;alert(1)&lt;/script&gt;" in html
    assert "&amp;" in html
    assert "&lt;img onerror=alert(1)&gt;" in html


def test_to_visual_wraps_builtin_html_as_fallback(tmp_path: Path):
    path = tmp_path / "visual.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "visual-fallback"
    prs.save(path)

    result = to_visual(path)

    assert result.kind == "pptx"
    assert result.fallback_html is not None
    assert "visual-fallback" in result.fallback_html


def test_to_visual_uses_fixed_safe_message_when_text_extract_fails(
    tmp_path: Path, monkeypatch
):
    path = tmp_path / "broken.pptx"
    path.write_bytes(b"x")

    def boom(_path: Path):
        raise RuntimeError(
            r"parse failed at C:\secret\dir\broken.pptx <script>alert(1)</script>"
        )

    monkeypatch.setattr(fmt_pptx, "to_html", boom)

    result = to_visual(path)

    assert result.kind == "pptx"
    assert result.fallback_html is not None
    assert "演示文稿已加密或损坏，无法生成文本回退。" in result.fallback_html
    assert "parse failed" not in result.fallback_html
    assert "C:\\secret\\dir\\broken.pptx" not in result.fallback_html
    assert "broken.pptx" not in result.fallback_html
    assert "<script>" not in result.fallback_html
