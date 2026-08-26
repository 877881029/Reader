from __future__ import annotations

import hashlib
import subprocess
import sys
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE


ROOT = Path(__file__).resolve().parents[1]
FIXTURE = ROOT / "tests" / "fixtures" / "pptx" / "visual-elements.pptx"
GENERATOR = ROOT / "scripts" / "generate_pptx_visual_fixture.py"
FIXTURE_SHA256 = "3ba6deda14de119b0de8751d5258461ea91f900634d7558c741ace3def96e8d4"


def test_fixture_generator_is_byte_deterministic_and_contains_real_elements(
    tmp_path,
):
    first = tmp_path / "first.pptx"
    second = tmp_path / "second.pptx"
    for output in (first, second):
        subprocess.run(
            [sys.executable, str(GENERATOR), "--output", str(output)],
            cwd=ROOT,
            check=True,
        )

    assert first.read_bytes() == second.read_bytes() == FIXTURE.read_bytes()
    assert hashlib.sha256(first.read_bytes()).hexdigest() == FIXTURE_SHA256

    with ZipFile(first) as archive:
        relationships = archive.read("ppt/_rels/presentation.xml.rels")
    assert relationships.index(b'/relationships/slide"') < relationships.index(
        b'/relationships/slideMaster"'
    )

    presentation = Presentation(first)
    assert len(presentation.slides) == 4
    assert (
        presentation.slides[0].background.fill.fore_color.rgb
        == RGBColor(20, 48, 90)
    )
    assert any(
        shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        for shape in presentation.slides[0].shapes
    )
    assert any(shape.has_table for shape in presentation.slides[1].shapes)
    assert any(shape.has_chart for shape in presentation.slides[2].shapes)
    assert (
        presentation.slides[3]
        .shapes[0]
        .text_frame.paragraphs[0]
        .runs[0]
        .font.name
        == "ReaderMissingFontZZ"
    )
