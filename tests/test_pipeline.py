from pathlib import Path

import pytest

from reader.preview.pipeline import preview
from reader.preview.result import PreviewResult
from reader.sniff import SniffError


class FakeOffice:
    def __init__(self, available=True, boom=False):
        self.available = available
        self.boom = boom
        self.calls = []
        self.available_calls = []

    def available_for(self, suffix: str) -> bool:
        self.available_calls.append(suffix)
        return self.available and suffix in {".docx", ".pptx", ".xlsx"}

    def export(self, path: Path) -> PreviewResult:
        self.calls.append(path)
        if self.boom:
            raise RuntimeError("COM busy")
        return PreviewResult(html="<p>office</p>", status_label="Office 预览", kind="html")


def test_markdown_default_is_visual_and_never_calls_office(tmp_path: Path):
    p = tmp_path / "a.md"
    p.write_text("# Z", encoding="utf-8")
    office = FakeOffice(available=True)
    result = preview(p, office=office)
    assert office.calls == []
    assert office.available_calls == []
    assert result.kind == "markdown"
    assert result.status_label == "内置预览（视觉模式）"
    assert result.fallback_html is not None
    assert "Z" in result.fallback_html


def test_markdown_visual_mode_supported_and_never_calls_office(tmp_path: Path):
    p = tmp_path / "visual.md"
    p.write_text("# Visual", encoding="utf-8")
    office = FakeOffice(available=True)

    result = preview(p, office=office, mode="visual")

    assert office.calls == []
    assert office.available_calls == []
    assert result.kind == "markdown"
    assert result.status_label == "内置预览（视觉模式）"
    assert result.fallback_html is not None
    assert "Visual" in result.fallback_html


def test_visual_mode_rejects_non_visual_suffix(tmp_path: Path):
    from docx import Document

    p = tmp_path / "a.docx"
    d = Document()
    d.add_paragraph("not visual")
    d.save(p)

    with pytest.raises(ValueError, match="visual mode supports only .pptx/.md"):
        preview(p, mode="visual")


def test_docx_defaults_to_builtin_without_office_call(tmp_path: Path):
    from docx import Document

    p = tmp_path / "a.docx"
    d = Document()
    d.add_paragraph("builtin-default")
    d.save(p)
    office = FakeOffice(available=True)
    result = preview(p, office=office)
    assert office.calls == []
    assert office.available_calls == []
    assert result.status_label == "内置预览"
    assert "builtin-default" in result.html


@pytest.mark.parametrize("suffix", [".docx", ".xlsx"])
def test_all_office_formats_are_builtin_first_without_any_com_probe(
    tmp_path: Path, monkeypatch, suffix: str
):
    from reader.preview import pipeline

    path = tmp_path / f"builtin-first{suffix}"
    path.write_bytes(b"fixture")
    office = FakeOffice(available=True)
    monkeypatch.setitem(
        pipeline._BUILTIN,
        suffix,
        lambda source: PreviewResult(
            html=f"<p>builtin-{source.suffix}</p>",
            status_label="内置预览",
        ),
    )

    result = preview(path, office=office)

    assert result.status_label == "内置预览"
    assert f"builtin-{suffix}" in result.html
    assert office.available_calls == []
    assert office.calls == []


def test_pptx_defaults_to_visual_without_any_com_probe(tmp_path: Path):
    from pptx import Presentation

    path = tmp_path / "builtin-first.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "visual-default"
    prs.save(path)
    office = FakeOffice(available=True)

    result = preview(path, office=office)

    assert result.kind == "pptx"
    assert result.status_label == "内置预览（视觉模式）"
    assert result.fallback_html is not None
    assert "visual-default" in result.fallback_html
    assert office.available_calls == []
    assert office.calls == []


def test_docx_explicit_office_mode_uses_export(tmp_path: Path):
    p = tmp_path / "a.docx"
    p.write_bytes(b"not-a-real-docx")
    office = FakeOffice(available=True)
    result = preview(p, office=office, mode="office")
    assert office.calls == [p]
    assert result.status_label == "Office 预览"
    assert "office" in result.html


def test_docx_falls_back_when_office_missing(tmp_path: Path):
    from docx import Document

    p = tmp_path / "a.docx"
    d = Document()
    d.add_paragraph("fallback-body")
    d.save(p)
    office = FakeOffice(available=False)
    result = preview(p, office=office, mode="office")
    assert office.calls == []
    assert result.status_label == "内置预览"
    assert "fallback-body" in result.html


def test_docx_propagates_when_export_raises(tmp_path: Path):
    p = tmp_path / "a.docx"
    p.write_bytes(b"not-a-real-docx")
    office = FakeOffice(available=True, boom=True)
    with pytest.raises(RuntimeError, match="COM busy"):
        preview(p, office=office, mode="office")


def test_pptx_uses_office_when_available(tmp_path: Path):
    p = tmp_path / "a.pptx"
    p.write_bytes(b"not-a-real-pptx")
    office = FakeOffice(available=True)
    result = preview(p, office=office, mode="office")
    assert office.calls == [p]
    assert result.status_label == "Office 预览"
    assert "office" in result.html


def test_xlsx_uses_office_when_available(tmp_path: Path):
    p = tmp_path / "a.xlsx"
    p.write_bytes(b"not-a-real-xlsx")
    office = FakeOffice(available=True)
    result = preview(p, office=office, mode="office")
    assert office.calls == [p]
    assert result.status_label == "Office 预览"
    assert "office" in result.html


def test_pptx_falls_back_when_office_missing(tmp_path: Path):
    from pptx import Presentation

    p = tmp_path / "a.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "pptx-fallback"
    prs.save(p)
    office = FakeOffice(available=False)
    result = preview(p, office=office, mode="office")
    assert office.available_calls == [".pptx"]
    assert office.calls == []
    assert result.kind == "pptx"
    assert result.fallback_html is not None
    assert "pptx-fallback" in result.fallback_html


def test_explicit_text_mode_returns_cacheable_html(tmp_path: Path):
    from pptx import Presentation

    path = tmp_path / "text.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "manual-text"
    prs.save(path)

    result = preview(path, mode="text")

    assert result.kind == "html"
    assert result.status_label == "内置预览（文本模式）"
    assert "manual-text" in result.html


def test_xlsx_falls_back_when_office_missing(tmp_path: Path):
    from openpyxl import Workbook

    p = tmp_path / "a.xlsx"
    wb = Workbook()
    wb.active.append(["xlsx-fallback-cell"])
    wb.save(p)
    office = FakeOffice(available=False)
    result = preview(p, office=office, mode="office")
    assert office.available_calls == [".xlsx"]
    assert office.calls == []
    assert result.status_label == "内置预览"
    assert "xlsx-fallback-cell" in result.html


def test_markdown_visual_renderer_error_propagates(tmp_path: Path, monkeypatch):
    from reader.preview import pipeline

    def boom(_path: Path) -> PreviewResult:
        raise ValueError("visual parse failed")

    monkeypatch.setattr(pipeline.fmt_md, "to_visual", boom)
    p = tmp_path / "a.md"
    p.write_text("# x", encoding="utf-8")
    with pytest.raises(ValueError, match="visual parse failed"):
        preview(p)


def test_unsupported_raises(tmp_path: Path):
    p = tmp_path / "x.pdf"
    p.write_bytes(b"%PDF")
    with pytest.raises(SniffError):
        preview(p)
